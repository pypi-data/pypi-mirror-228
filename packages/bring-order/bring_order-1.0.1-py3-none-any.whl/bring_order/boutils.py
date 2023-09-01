"""Helpful Javascript methods"""
import textwrap
from pptx.util import Inches, Pt
from IPython.display import display, Javascript


class BOUtils:
    """Helpful Javascript methods"""
    def __init__(self):
        """Class constructor.
        The instance attribute self.change_cell_count is a function that
        can be manipulated to change the cell_count attribute of the Bodi/Deductive/Inductive
        instance that is calling BOUtils methods.
        """

        self.change_cell_count = lambda n: None
        self.prs = None
        self.pptx_file = ''

    def create_code_cells_above(self, how_many):
        """Creates code cells above the current cell
        
        Args:
            how_many (int): the number of cells to be created
        """

        for _ in range(how_many):
            command = '''
            IPython.notebook.insert_cell_above("code")
            '''
            display(Javascript(command))

    def focus_on_input_above(self, distance):
        ''' Sets cursor focus on input above.'''
        command = f'''
            var cell = IPython.notebook.get_cell(-{distance})
            cell.code_mirror.getInputField().focus()
        '''
        display(Javascript(command))

    def create_markdown_cells_above(self, how_many, text='', run_cell=True):
        """Creates markdown cells above the current cell
        
        Args:
            how_many (int): the number of cells to be created
            text (str): default contents of the created cells
            run_cell (bool): should the cell be run after creation
        """

        if run_cell:
            run_cell_str = 'true'
        else:
            run_cell_str = 'false'

        for _ in range(how_many):
            command = f'''
                var cell = IPython.notebook.insert_cell_above("markdown");
                cell.set_text("{text}");
                if ({run_cell_str}) {{
                    var new_index = IPython.notebook.get_selected_index()-1;
                    IPython.notebook.execute_cells([new_index]);
                }}
            '''
            display(Javascript(command))
        self.__create_slide(text)

    def __create_slide(self, context):
        ''' Creates new slide to slide deck
            Args:
                context: str
        '''

        rows = context.split('\\n')
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        left = top = Inches(1)
        width = height = Inches(6)
        txt_box = slide.shapes.add_textbox(left, top, width, height)
        txt_frame = txt_box.text_frame

        for row in rows:
            if '#' in row:
                parag = txt_frame.add_paragraph()
                row = row.replace('#','')
                row = textwrap.fill(row, 85)
                parag.text = row
                parag.font.size = Pt(22)
                parag.font.bold = True
            else:
                parag = txt_frame.add_paragraph()
                row = row.replace('               ','')
                row = textwrap.fill(row, 85)
                parag.text = row
                parag.font.size = Pt(16)
        del slide
        self.prs.save(self.pptx_file)

    def clear_code_cells_above(self, how_many):
        """Clears code cells above the active cell
        
        Args:
            how_many (int): the number of cells to be cleared
        """

        for _ in range(how_many):
            self.delete_cell_above()
        self.create_code_cells_above(how_many)

    def delete_cell_above(self):
        """Deletes the cell above the current cell"""

        command = '''
        var above_index = IPython.notebook.get_selected_index() - 1;
        IPython.notebook.delete_cell(above_index);
        '''
        display(Javascript(command))

    def delete_cell_from_current(self, distance):
        """Deletes a cell that has the index of the active cell index + distance
        
        Args:
            distance (int): target cell index with respect to active cell
        """

        command = f'''
        var index = IPython.notebook.get_selected_index() + {distance};
        IPython.notebook.delete_cell(index);
        '''
        display(Javascript(command))

    def run_cells_above(self, cell_count):
        """Runs cells above the active cell.
        
        Args:
            cell_count (int): the number of cells to be run
        """

        command = f'''
        var output_area = this;
        var cell_element = output_area.element.parents('.cell');
        var current_index = Jupyter.notebook.get_cell_elements().index(cell_element);
        var first_index = current_index - {cell_count};
        IPython.notebook.execute_cell_range(first_index, current_index);
        '''

        display(Javascript(command))

    def execute_cell_from_current(self, distance, code='', hide_input=True):
        """Executes code in cell that has the index of the active cell index + distance
        
        Args:
            distance (int): target cell index with respect to active cell
            code (str): Python code to be executed
            hide_input (boolean): hides input of the executed cell, defaults to True 
        """

        if hide_input:
            hide_input_string = 'true'
        else:
            hide_input_string = 'false'

        command = f'''
        var index = IPython.notebook.get_selected_index() + {distance};
        var cell = IPython.notebook.get_cell(index);
        cell.set_text("{code}");
        IPython.notebook.execute_cells([index]);
        if ({hide_input_string}) (cell.input.hide());
        '''

        display(Javascript(command))

    def hide_selected_input(self):
        """Hides the input of the selected cell"""

        command = '''
        var cell_index = IPython.notebook.get_selected_index();
        var cells = IPython.notebook.get_cells();
        cells[cell_index].input.hide();
        '''

        display(Javascript(command))

    def hide_current_input(self):
        """Hides the input of the currently executing cell"""

        command = '''
        var output_area = this;
        var cell_element = output_area.element.parents('.cell');
        var cell_idx = Jupyter.notebook.get_cell_elements().index(cell_element);
        var cell = Jupyter.notebook.get_cell(cell_idx);
        cell.input.hide();
        '''

        display(Javascript(command))

    def print_to_console(self, msg):
        """Prints to browser console. Useful for debugging etc.
        
        Args:
            msg (str)
        """

        command = f'''
        console.log("{msg}");
        '''

        display(Javascript(command))

    def check_cells_above(self, cell_count, test_name, variables):
        """Check if cells above contain the given test and at least one of the given variables.
        Prints a warning.

        Args:
            cell_count (int): the number of cells to be checked
            test_name (str): the text you are trying to find
            variables (list): list of strings
        """

        command = f'''
        var current_index = IPython.notebook.get_selected_index();
        var first_index = current_index - {cell_count};
        var cells = IPython.notebook.get_cells();
        var variables = {variables}
        var warning = ""
        cells.forEach(function(cell) {{
                var index = IPython.notebook.find_cell_index(cell);
                if(index >= first_index && index < current_index) {{
                    var cell_text = cell.get_text();
                    if(cell_text.includes("{test_name}")) {{
                        variables.forEach(function(variable) {{
                            if(cell_text.includes(variable)) {{
                                warning = warning + " " + variable + ",";
                            }}
                        }});
                    }}
                }}
        }});
        if(warning != "") {{
            warning = "Warning! It seems that you are trying to use " + "{test_name}" + " for variables that are not normally distributed:" + warning
            warning = warning.slice(0, -1);
            element.text(warning);
        }}
        '''

        display(Javascript(command))

    def get_python_code_from_response(self, response):
        """Returns the Python code extracted from AI response.
        
        Args:
            response (str): the AI response
            
        Returns:
            code (str): the Python code from the response
        """

        substrings = []
        start = 0
        while True:
            start = response.find('```python\n', start)
            if start == -1:
                break
            start += len('```python"')
            end = response.find('\n```', start)
            if end == -1:
                break
            code = response[start:end]
            code = code.replace('\\', '\\\\')
            code = code.replace('\n', '\\n')
            code = code.replace('"', '\\"')
            code = code.replace("'", "\\'")
            substrings.append(code)
            start = end

        if substrings:
            self.print_to_console("\\n".join(substrings))
            return "\\n".join(substrings)

        return 'No Python code in the response'

    def insert_ai_code_into_new_cell(self, code=''):
        """Opens new empty code cell and inserts the given code into it.
        
        Args:
            code (str)
        """

        command1 = '''
        IPython.notebook.insert_cell_above("code");
        var index = IPython.notebook.get_selected_index() - 1;
        var cell = IPython.notebook.get_cell(index);
        var cell_above = IPython.notebook.get_cell(index - 1);
        while (cell_above.get_text() == "") {
            cell = cell_above;
            index = index - 1;
            cell_above = IPython.notebook.get_cell(index - 1);
        }
        cell.set_text("# If the code doesn't appear here, please, copy-paste it manually.");
        '''

        display(Javascript(command1))
        self.change_cell_count(1)

        command2 = f'''
        var output_area = this;
        var active_element = output_area.element.parents('.cell');
        var index = Jupyter.notebook.get_cell_elements().index(active_element) - 1;
        var cell = Jupyter.notebook.get_cell(index);
        while (cell.get_text() != "# If the code doesn't appear here, please, copy-paste it manually.") {{
            index = index - 1;
            cell = IPython.notebook.get_cell(index);
        }}
        cell.set_text("{code}");
        '''

        display(Javascript(command2))
